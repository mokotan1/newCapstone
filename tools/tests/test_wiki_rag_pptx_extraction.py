from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from wiki_rag.extractors.pptx import extract_pptx


@pytest.fixture
def pptx_fixture(tmp_path: Path) -> Path:
    presentation = Presentation()
    first_slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    first_text = first_slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    )
    first_text.text = "첫 번째 슬라이드"
    table_shape = first_slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(2),
        Inches(4),
        Inches(1),
    )
    table_shape.table.cell(0, 0).text = "열 A"
    table_shape.table.cell(0, 1).text = "열 B"
    table_shape.table.cell(1, 0).text = "값 1"
    table_shape.table.cell(1, 1).text = "값 2"
    image_path = tmp_path / "visual.png"
    Image.new("RGB", (2, 2), color="red").save(image_path)
    first_slide.shapes.add_picture(
        str(image_path),
        Inches(5),
        Inches(1),
        Inches(1),
        Inches(1),
    )
    first_slide.notes_slide.notes_text_frame.text = "발표자 메모"

    second_slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    second_text = second_slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    )
    second_text.text = "두 번째 슬라이드"

    fixture = tmp_path / "발표 자료.pptx"
    presentation.save(fixture)
    return fixture


def test_pptx_markdown_keeps_slide_order_and_speaker_notes(
    pptx_fixture: Path,
) -> None:
    result = extract_pptx(pptx_fixture)

    assert "## Slide 1" in result.markdown
    assert "발표자 메모" in result.markdown
    assert result.page_or_slide_count == 2
    assert result.markdown.index("첫 번째 슬라이드") < result.markdown.index(
        "두 번째 슬라이드"
    )


def test_pptx_tables_are_rendered_as_markdown(pptx_fixture: Path) -> None:
    result = extract_pptx(pptx_fixture)

    assert "| 열 A | 열 B |" in result.markdown
    assert "| --- | --- |" in result.markdown
    assert "| 값 1 | 값 2 |" in result.markdown


def test_pptx_visual_assets_emit_inspection_placeholder(
    pptx_fixture: Path,
) -> None:
    result = extract_pptx(pptx_fixture)

    assert (
        "> Visual asset present: slide 1, shape 3; "
        "inspect original PPTX/PDF."
    ) in result.markdown
    assert "pptx_visual_asset:slide_1:shape_3" in result.warnings
