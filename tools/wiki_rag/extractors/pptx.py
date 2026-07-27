"""Extract ordered text, tables, notes, and visual signals from PPTX files."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from . import ExtractionResult

_VISUAL_SHAPE_TYPES = frozenset(
    {
        MSO_SHAPE_TYPE.CHART,
        MSO_SHAPE_TYPE.DIAGRAM,
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.IGX_GRAPHIC,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
        MSO_SHAPE_TYPE.MEDIA,
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.WEB_VIDEO,
    }
)


def _normalize_text(value: str) -> str:
    return value.replace("\r\n", "\n").replace("\r", "\n").strip()


def _escape_table_cell(value: str) -> str:
    normalized = _normalize_text(value).replace("|", r"\|")
    return normalized.replace("\n", "<br>")


def _render_table(table: Any) -> list[str]:
    rows = [
        [_escape_table_cell(cell.text) for cell in row.cells]
        for row in table.rows
    ]
    if not rows:
        return []

    width = len(rows[0])
    lines = [
        "| " + " | ".join(rows[0]) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in rows[1:])
    return lines


def extract_pptx(source_path: Path) -> ExtractionResult:
    """Extract slides in document shape order, including speaker notes."""

    presentation = Presentation(source_path)
    sections: list[str] = []
    warnings: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        sections.extend((f"## Slide {slide_number}", ""))
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_table", False):
                table_lines = _render_table(shape.table)
                if table_lines:
                    sections.extend((*table_lines, ""))
                continue

            text = ""
            if getattr(shape, "has_text_frame", False):
                text = _normalize_text(shape.text)
            if text:
                sections.extend((text, ""))
                continue

            if shape.shape_type in _VISUAL_SHAPE_TYPES:
                sections.extend(
                    (
                        (
                            "> Visual asset present: "
                            f"slide {slide_number}, shape {shape_number}; "
                            "inspect original PPTX/PDF."
                        ),
                        "",
                    )
                )
                warnings.append(
                    "pptx_visual_asset:"
                    f"slide_{slide_number}:shape_{shape_number}"
                )

        notes = _normalize_text(slide.notes_slide.notes_text_frame.text)
        if notes:
            sections.extend(("### Speaker notes", "", notes, ""))

    return ExtractionResult(
        markdown="\n".join(sections).rstrip() + "\n",
        page_or_slide_count=len(presentation.slides),
        warnings=warnings,
    )
